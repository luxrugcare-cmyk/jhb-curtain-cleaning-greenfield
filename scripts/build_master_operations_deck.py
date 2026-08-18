#!/usr/bin/env python3
"""Builds a high-impact, professional widescreen PowerPoint presentation (.pptx)
for the JHB Curtain Cleaning Master Systems Operations & Onboarding Manual.
"""

import os
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand Palette (Craftsman Authority / Metallic Champagne)
DARK_BG = RGBColor(18, 18, 18)       # #121212
CARD_BG = RGBColor(26, 26, 26)       # #1A1A1A
GOLD = RGBColor(201, 156, 45)         # #C99C2D
LIGHT_GOLD = RGBColor(226, 190, 88)   # #E2BE58
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(212, 212, 212)
MUTED = RGBColor(163, 163, 163)
GREEN = RGBColor(52, 211, 153)

def create_slide_deck():
    prs = Presentation()
    # 16:9 Widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide_bg(slide):
        # Full dark background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category="JHB CURTAIN CLEANING · MASTER OPERATIONS MANUAL"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        p = cat_tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = GOLD
        
        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        t_tf = t_box.text_frame
        t_tf.word_wrap = True
        p2 = t_tf.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

        # Gold accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.fill.background()

    # -------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s1)

    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "JHB CURTAIN CLEANING"
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = LIGHT_GOLD
    p0.space_after = Pt(8)

    p1 = tf.add_paragraph()
    p1.text = "Master Systems Operations & Onboarding Manual"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "Complete Architecture, Autonomous Outreach, Meta & WhatsApp Cloud API, Attio CRM, and Daily SOPs"
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_GRAY
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "Lead Operator: Stephen (+27 75 011 9200) · Platform Release: v0.5.0 Enterprise · jhbcurtaincleaning.co.za"
    p3.font.size = Pt(12)
    p3.font.color.rgb = GOLD

    # -------------------------------------------------------------
    # SLIDE 2: 4-TIER ARCHITECTURE & DATA FLOW
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s2)
    add_header(s2, "System Architecture & 4-Tier Data Pipeline")

    cards = [
        ("1. Inbound Channels", [
            "• Next.js Web Portal (48 Canonical Routes)",
            "• Meta Ads (FB/IG Click-to-WhatsApp)",
            "• WhatsApp Direct (+27 75 011 9200)",
            "• B2B Cold Outreach (800 Scored Leads)",
        ], Inches(0.8)),
        ("2. Ingestion Engine", [
            "• /api/leads API Route",
            "• /api/webhooks/whatsapp",
            "• Vercel Blob (3 Photo Attachments)",
            "• 4-Tier Fallback Error Recovery",
        ], Inches(3.8)),
        ("3. Core Operations", [
            "• Attio CRM (1,705 Clean Contacts)",
            "• AgentMail (stephen-1015@agentmail.to)",
            "• Meta Conversions API (SHA-256 CAPI)",
            "• Sanity Studio CMS (/studio)",
        ], Inches(6.8)),
        ("4. Automated Responses", [
            "• Residential Welcome & Zero Shrinkage",
            "• Commercial Protocol & SANS 1423",
            "• 10% Trade Partner Kit & Referral Code",
            "• 5-Star Google Review Solicitation",
        ], Inches(9.8)),
    ]

    for title, bullet_points, left_pos in cards:
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(2.7), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = GOLD
        card.line.width = Pt(1.5)

        tb = s2.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.9), Inches(2.4), Inches(4.8))
        c_tf = tb.text_frame
        c_tf.word_wrap = True
        
        p = c_tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = LIGHT_GOLD
        p.space_after = Pt(12)

        for pt in bullet_points:
            bp = c_tf.add_paragraph()
            bp.text = pt
            bp.font.size = Pt(11)
            bp.font.color.rgb = LIGHT_GRAY
            bp.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 3: LEAD CAPTURE & CONVERSION ENGINE
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s3)
    add_header(s3, "Lead Capture & Ingestion Engine")

    intake_points = [
        ("Residential Assessment Wizard (/quote)", [
            "• Dynamic quote builder with fabric & room selector",
            "• Photo uploader (up to 3 photos saved to Vercel Blob)",
            "• Instant Welcome Email with 'Do Not Take Down' notice",
            "• Direct WhatsApp escalation button for urgent bookings",
        ], Inches(0.8), Inches(1.8)),
        ("Commercial Assessment Portal (/commercial-assessment)", [
            "• Hotel, executive lodge, and corporate facility intake",
            "• Room count, sector selection, and timeline tracking",
            "• Automatic dispatch of SANS 1423 Fire Retardancy Protocol",
            "• Zero Room Downtime guarantee (10:00–14:00 guest turnover)",
        ], Inches(6.8), Inches(1.8)),
        ("10% Trade Partner Hub (/trade)", [
            "• Dedicated portal for Interior Decorators & Curtain Workrooms",
            "• Interactive Passive Revenue Earnings Calculator",
            "• Generates unique Studio Referral Code (e.g. TRADE-HYDEPA-4821)",
            "• 48-Hour EFT commission payout tracking",
        ], Inches(0.8), Inches(4.5)),
        ("Multi-Channel Lead Resilience Engine", [
            "• Attio CRM ingestion with UUID idempotency key",
            "• AgentMail notifications to stephen-1015@agentmail.to",
            "• Meta Conversions API (CAPI) server-side event tracking",
            "• Local fallback archive ensures 0% lead loss during outages",
        ], Inches(6.8), Inches(4.5)),
    ]

    for title, bullets, left_pos, top_pos in intake_points:
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.7), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = GOLD

        tb = s3.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.1), Inches(5.4), Inches(2.2))
        c_tf = tb.text_frame
        c_tf.word_wrap = True

        p = c_tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = LIGHT_GOLD
        p.space_after = Pt(6)

        for pt in bullets:
            bp = c_tf.add_paragraph()
            bp.text = pt
            bp.font.size = Pt(10.5)
            bp.font.color.rgb = LIGHT_GRAY
            bp.space_after = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 4: AUTONOMOUS B2B EMAIL OUTREACH ENGINE
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s4)
    add_header(s4, "Autonomous B2B Email Outreach Engine")

    touches = [
        ("Touch 1 (Day 1): Plain-Text-First Peer Inquiry", [
            "• Clean, high-deliverability plain text note from Stephen",
            "• Zero tracking pixels or heavy HTML to bypass spam filters",
            "• Short, direct question tailored to Hotel GMs or Designers",
            "• Focus: Quick response & opening a conversation",
        ], Inches(0.8)),
        ("Touch 2 (Day 4): High-Visual Metallic MJML Proof Card", [
            "• Google Stitch metallic styling with The Leonardo case study",
            "• SANS 1423 flame retardancy certification badges",
            "• Zero room downtime workflow breakdown (10:00–14:00)",
            "• One-click WhatsApp link to Stephen (+27 75 011 9200)",
        ], Inches(4.8)),
        ("Touch 3 (Day 8): Graceful Breakup Note", [
            "• Short, respectful 2-sentence closing note",
            "• Leaves door open for seasonal Highveld dust storms",
            "• Provides direct WhatsApp number for future reference",
            "• 100/100 Anti-Spam Deliverability audit passed",
        ], Inches(8.8)),
    ]

    for title, bullets, left_pos in touches:
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(3.7), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = GOLD

        tb = s4.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.9), Inches(3.4), Inches(4.8))
        c_tf = tb.text_frame
        c_tf.word_wrap = True

        p = c_tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = LIGHT_GOLD
        p.space_after = Pt(10)

        for pt in bullets:
            bp = c_tf.add_paragraph()
            bp.text = pt
            bp.font.size = Pt(11)
            bp.font.color.rgb = LIGHT_GRAY
            bp.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 5: META (FB/IG) & WHATSAPP BUSINESS CLOUD API
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s5)
    add_header(s5, "Meta Ads & WhatsApp Business Cloud API")

    meta_items = [
        ("WhatsApp Webhook Endpoint (/api/webhooks/whatsapp)", [
            "• Callback URL: https://www.jhbcurtaincleaning.co.za/api/webhooks/whatsapp",
            "• Verify Token: jhb_curtain_cleaning_meta_2026",
            "• Automated GET hub.challenge verification with HTTP 200",
            "• Ingests incoming customer messages & curtain photos directly",
        ], Inches(0.8), Inches(1.8)),
        ("3-Button Interactive WhatsApp Menu", [
            "• Dispatches instant interactive menu on first customer message",
            "• Button 1: Residential Clean (On-site zero shrinkage)",
            "• Button 2: Commercial / Hotel (Zero room downtime)",
            "• Button 3: 10% Trade Partner (Decorator & workroom)",
        ], Inches(6.8), Inches(1.8)),
        ("Meta Conversions API (CAPI) & Pixel", [
            "• Server-side CAPI tracking for Facebook & Instagram Ads",
            "• SHA-256 hashing for all user data (Email, Phone, City)",
            "• Fires 'Lead' event on every quote & commercial assessment",
            "• Verified Meta Business ID: 1164022258518231 | Asset: 867076983156221",
        ], Inches(0.8), Inches(4.5)),
        ("Click-to-WhatsApp (CTWA) Ad Campaigns", [
            "• Luxury Residential Ad Set (Sandhurst, Hyde Park, Pretoria East)",
            "• Hotels & Lodges Ad Set (Zero lost room nights)",
            "• Trade Decorators Ad Set (10% passive referral revenue)",
            "• Highveld Dust Storm Seasonal Retargeting Campaign",
        ], Inches(6.8), Inches(4.5)),
    ]

    for title, bullets, left_pos, top_pos in meta_items:
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.7), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = GOLD

        tb = s5.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.1), Inches(5.4), Inches(2.2))
        c_tf = tb.text_frame
        c_tf.word_wrap = True

        p = c_tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = LIGHT_GOLD
        p.space_after = Pt(6)

        for pt in bullets:
            bp = c_tf.add_paragraph()
            bp.text = pt
            bp.font.size = Pt(10.5)
            bp.font.color.rgb = LIGHT_GRAY
            bp.space_after = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 6: STEPHEN'S DAILY OPERATIONAL SOP
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s6)
    add_header(s6, "Operator Daily Standard Operating Procedure (SOP)")

    schedule = [
        ("07:30 AM · Inbound Review", "Check AgentMail (stephen-1015@agentmail.to) for overnight quote requests & photos.", Inches(0.8), Inches(1.8)),
        ("08:00 AM · WhatsApp Schedule", "Confirm on-site arrival times with scheduled residential clients via WhatsApp.", Inches(0.8), Inches(2.8)),
        ("09:00 AM · B2B Outreach", "Run daily warm-up batch: python scripts/b2b_multi_sequence_dispatcher.py --segment hotel", Inches(0.8), Inches(3.8)),
        ("04:00 PM · Proposal Generation", "Generate commercial quotations for site inspections using generate_commercial_proposal.py.", Inches(0.8), Inches(4.8)),
        ("05:30 PM · 5-Star Reviews", "Trigger automated Google Review review request links for completed jobs.", Inches(0.8), Inches(5.8)),
    ]

    for time_title, desc, left_pos, top_pos in schedule:
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = GOLD

        tb = s6.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.08), Inches(11.3), Inches(0.7))
        c_tf = tb.text_frame
        c_tf.word_wrap = True

        p = c_tf.paragraphs[0]
        p.text = time_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = LIGHT_GOLD

        p2 = c_tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT_GRAY

    # -------------------------------------------------------------
    # SLIDE 7: COMMAND CHEAT SHEET & VERIFICATION
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_bg(s7)
    add_header(s7, "Terminal Maintenance & Execution Commands")

    cmd_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    cmd_box.fill.solid()
    cmd_box.fill.fore_color.rgb = CARD_BG
    cmd_box.line.color.rgb = GOLD

    tb = s7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.6))
    c_tf = tb.text_frame
    c_tf.word_wrap = True

    commands = [
        ("Run Full Test Suite (6 Suites):", "npm run test:all"),
        ("Test All Lifecycle Automation Triggers:", "npm run test:triggers"),
        ("Test Meta CAPI & WhatsApp Webhook:", "npm run test:meta"),
        ("Crawl & Verify All 48 Routes:", "npm run test:crawl"),
        ("Audit Email Spam Deliverability (19 Templates):", "python scripts/audit_spam_score.py"),
        ("Generate Commercial PDF Quotation:", "python scripts/generate_commercial_proposal.py --client 'Name' --company 'Hotel'"),
        ("Submit Sitemap to Google Search Console:", "python scripts/submit_gsc_indexing.py"),
    ]

    p = c_tf.paragraphs[0]
    p.text = "OPERATIONAL TERMINAL COMMANDS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = LIGHT_GOLD
    p.space_after = Pt(10)

    for label, cmd in commands:
        p_lbl = c_tf.add_paragraph()
        p_lbl.text = f"{label}"
        p_lbl.font.size = Pt(11)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = WHITE
        
        p_cmd = c_tf.add_paragraph()
        p_cmd.text = f"  $ {cmd}"
        p_cmd.font.size = Pt(11)
        p_cmd.font.color.rgb = GREEN
        p_cmd.space_after = Pt(6)

    # Save to public downloads and docs
    out_dir_public = pathlib.Path("public/downloads")
    out_dir_public.mkdir(parents=True, exist_ok=True)
    out_file_public = out_dir_public / "JHB_Curtain_Cleaning_Master_Operations_Manual.pptx"
    prs.save(str(out_file_public))

    out_dir_docs = pathlib.Path("docs")
    out_file_docs = out_dir_docs / "JHB_Curtain_Cleaning_Master_Operations_Manual.pptx"
    prs.save(str(out_file_docs))

    print(f"✓ PowerPoint slide deck created successfully: {out_file_public}")
    print(f"✓ PowerPoint slide deck copied to: {out_file_docs}")

if __name__ == "__main__":
    create_slide_deck()
